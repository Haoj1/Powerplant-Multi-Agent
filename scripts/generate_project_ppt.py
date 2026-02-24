#!/usr/bin/env python3
"""
Generate a ~30-slide PowerPoint presentation for the Multi-Agent Powerplant Monitoring System.
Run: pip install python-pptx matplotlib
      python scripts/generate_project_ppt.py
Output: docs/Multi-Agent_Powerplant_Project.pptx
"""

import os
import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError as e:
    print(f"ImportError: {e}")
    print("Please install: pip install python-pptx")
    sys.exit(1)

try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
    from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
    HAS_MATPLOTLIB = True
except ImportError:
    HAS_MATPLOTLIB = False
    print("matplotlib not found; diagrams will be placeholders. Install: pip install matplotlib")


def create_architecture_diagram(output_path: Path) -> bool:
    """Create architecture diagram and save as PNG."""
    if not HAS_MATPLOTLIB:
        return False
    fig, ax = plt.subplots(1, 1, figsize=(10, 6))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 6)
    ax.axis("off")

    # Colors
    sim_color = "#4CAF50"
    agent_colors = ["#2196F3", "#FF9800", "#9C27B0", "#E91E63"]
    mqtt_color = "#607D8B"

    # MQTT bus (center)
    ax.add_patch(FancyBboxPatch((1, 2.2), 8, 1.6, boxstyle="round,pad=0.05",
                                facecolor=mqtt_color, edgecolor="white", alpha=0.8))
    ax.text(5, 3, "MQTT Message Bus\n(telemetry / alerts / diagnosis / tickets)", ha="center", va="center",
            fontsize=10, color="white", fontweight="bold")

    # Simulator (top)
    ax.add_patch(FancyBboxPatch((3.5, 4.5), 3, 0.8, boxstyle="round,pad=0.05",
                                facecolor=sim_color, edgecolor="white", alpha=0.9))
    ax.text(5, 4.9, "Simulator", ha="center", va="center", fontsize=11, fontweight="bold")
    ax.annotate("", xy=(5, 3.8), xytext=(5, 4.5), arrowprops=dict(arrowstyle="->", color="gray", lw=2))

    # Agents (bottom)
    agents = ["Agent A\n(Monitor)", "Agent B\n(Diagnosis)", "Agent C\n(Ticket)", "Agent D\n(Review)"]
    for i, (name, color) in enumerate(zip(agents, agent_colors)):
        x = 1.2 + i * 2.2
        ax.add_patch(FancyBboxPatch((x, 0.3), 1.8, 1.2, boxstyle="round,pad=0.05",
                                    facecolor=color, edgecolor="white", alpha=0.9))
        ax.text(x + 0.9, 0.9, name, ha="center", va="center", fontsize=8, color="white", fontweight="bold")
        ax.annotate("", xy=(x + 0.9, 2.2), xytext=(x + 0.9, 1.5),
                   arrowprops=dict(arrowstyle="->", color="gray", lw=1.5))

    ax.set_title("System Architecture", fontsize=14, fontweight="bold")
    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()
    return True


def create_data_flow_diagram(output_path: Path) -> bool:
    """Create data flow diagram."""
    if not HAS_MATPLOTLIB:
        return False
    fig, ax = plt.subplots(1, 1, figsize=(10, 5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 5)
    ax.axis("off")

    nodes = [
        (1, 2.5, "Telemetry\n(sensors)", "#4CAF50"),
        (3, 2.5, "Alerts", "#FF9800"),
        (5, 2.5, "Diagnosis\n(RCA)", "#9C27B0"),
        (7, 2.5, "Ticket", "#E91E63"),
        (9, 2.5, "Feedback", "#00BCD4"),
    ]
    for x, y, label, color in nodes:
        ax.add_patch(FancyBboxPatch((x - 0.5, y - 0.4), 1, 0.8, boxstyle="round,pad=0.05",
                                    facecolor=color, edgecolor="white", alpha=0.9))
        ax.text(x, y, label, ha="center", va="center", fontsize=9, color="white", fontweight="bold")

    for i in range(len(nodes) - 1):
        ax.annotate("", xy=(nodes[i + 1][0] - 0.5, nodes[i + 1][1]),
                   xytext=(nodes[i][0] + 0.5, nodes[i][1]),
                   arrowprops=dict(arrowstyle="->", color="gray", lw=2))

    ax.text(5, 4.2, "Data Flow: Simulator → Monitor → Diagnosis → Ticket → Review", ha="center", fontsize=11)
    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()
    return True


def create_dashboard_mockup(output_path: Path) -> bool:
    """Create simple dashboard layout mockup."""
    if not HAS_MATPLOTLIB:
        return False
    fig, ax = plt.subplots(1, 1, figsize=(8, 5))
    ax.set_xlim(0, 8)
    ax.set_ylim(0, 5)
    ax.axis("off")

    # Header
    ax.add_patch(FancyBboxPatch((0, 4.2), 8, 0.7, boxstyle="round,pad=0.02", facecolor="#1976D2", alpha=0.9))
    ax.text(4, 4.55, "Agent D Review Dashboard", ha="center", va="center", fontsize=14, color="white", fontweight="bold")

    # Sidebar
    ax.add_patch(FancyBboxPatch((0, 0), 1.5, 4.2, boxstyle="round,pad=0.02", facecolor="#37474F", alpha=0.8))
    for i, label in enumerate(["Review Queue", "Alerts", "Sensors", "Chat", "Scenarios"]):
        ax.text(0.75, 3.8 - i * 0.7, label, ha="center", va="center", fontsize=9, color="white")

    # Main content
    ax.add_patch(FancyBboxPatch((1.6, 0.2), 6.3, 3.9, boxstyle="round,pad=0.02", facecolor="#ECEFF1", alpha=0.9))
    ax.text(4.75, 2.15, "Review Queue\nDiagnosis List\nApprove / Reject", ha="center", va="center", fontsize=11)

    ax.set_title("Dashboard UI", fontsize=12, fontweight="bold")
    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close()
    return True


def add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    tf = slide.shapes.add_textbox(left, top, width, Inches(1.5))
    p = tf.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        tf2 = slide.shapes.add_textbox(left, Inches(3.8), width, Inches(1))
        p2 = tf2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.alignment = PP_ALIGN.CENTER


def add_content_slide(prs: Presentation, title: str, bullets: list, image_path: Path = None):
    """Add a content slide with optional image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    p = tf.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(5.5), Inches(1.2), width=Inches(4), height=Inches(4.5))
        content_width = Inches(4.8)
    else:
        content_width = Inches(9)

    # Bullets
    tf2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), content_width, Inches(5.5))
    text_frame = tf2.text_frame
    text_frame.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.space_after = Pt(12)


def main():
    output_pptx = PROJECT_ROOT / "docs" / "Multi-Agent_Powerplant_Project.pptx"
    output_pptx.parent.mkdir(parents=True, exist_ok=True)

    img_dir = PROJECT_ROOT / "docs" / "ppt_images"
    img_dir.mkdir(parents=True, exist_ok=True)

    # Generate diagrams
    arch_path = img_dir / "architecture.png"
    flow_path = img_dir / "data_flow.png"
    dash_path = img_dir / "dashboard.png"
    create_architecture_diagram(arch_path)
    create_data_flow_diagram(flow_path)
    create_dashboard_mockup(dash_path)

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slides = [
        ("Title", add_title_slide, "Multi-Agent Powerplant Monitoring System",
         "Real-time Anomaly Detection, Root Cause Analysis & Human-in-the-Loop Review"),
        ("Outline", add_content_slide, "Presentation Outline", [
            "1. Why This Project?",
            "2. Importance & Motivation",
            "3. Problems We Solve",
            "4. System Architecture",
            "5. Agent Roles & Data Flow",
            "6. Key Features & Demo",
            "7. Technology Stack",
            "8. Conclusion & Future Work",
        ]),
        ("Why This Project", add_content_slide, "Why We Chose This Project", [
            "• Combines multiple cutting-edge areas: Multi-Agent Systems, IoT, LLM/AI",
            "• Practical industrial application: powerplant asset monitoring",
            "• End-to-end pipeline: from raw sensors to human-approved tickets",
            "• Human-in-the-loop design ensures safety and accountability",
            "• Extensible: rules, RAG, Salesforce integration, multimodal (3D + VLM)",
        ]),
        ("Importance", add_content_slide, "Importance & Motivation", [
            "• Industrial IoT generates massive sensor data; manual monitoring is costly",
            "• Early anomaly detection prevents equipment failure and downtime",
            "• Root cause analysis (RCA) reduces mean-time-to-repair (MTTR)",
            "• Automated ticket creation streamlines maintenance workflows",
            "• Human review loop maintains quality and builds trust",
        ]),
        ("Problems", add_content_slide, "Problems We Solve", [
            "• Real-time anomaly detection from streaming telemetry",
            "• Interpretable root cause diagnosis (not black-box)",
            "• Automated ticket creation with evidence and recommended actions",
            "• Human review interface with approve/reject/edit",
            "• Integration with external systems (e.g., Salesforce)",
        ]),
        ("Challenges", add_content_slide, "Key Challenges", [
            "• Low latency: detect and diagnose within seconds",
            "• Explainability: evidence and rules, not just predictions",
            "• Scalability: multiple assets, MQTT pub/sub",
            "• Reproducibility: fault injection, scenario-driven testing",
        ]),
        ("Domain", add_content_slide, "Problem Domain: Industrial IoT", [
            "• Powerplants, pumps, bearings, valves generate continuous sensor data",
            "• Pressure, flow, temperature, vibration, motor current, etc.",
            "• Faults: bearing wear, clogging, valve stuck, sensor drift",
            "• Early detection prevents costly downtime and safety incidents",
        ]),
        ("Approach", add_content_slide, "Our Approach: Multi-Agent Pipeline", [
            "• 4 specialized agents + Simulator, coordinated via MQTT",
            "• Agent A: Threshold / Z-score anomaly detection",
            "• Agent B: Rule-based + LLM (ReAct) diagnosis",
            "• Agent C: Ticket creation (GitHub Issues / local)",
            "• Agent D: Human review, feedback, Salesforce Case creation",
        ]),
        ("Architecture", add_content_slide, "System Architecture", [
            "• Simulator: Generates telemetry, injects faults (bearing_wear, clogging, etc.)",
            "• MQTT: Message bus for telemetry, alerts, diagnosis, tickets",
            "• Agents: Subscribe/publish, stateless, scalable",
            "• Agent D Dashboard: React frontend for review, chat, scenarios",
        ], arch_path),
        ("Architecture Diagram", add_content_slide, "Architecture Diagram", [
            "Simulator publishes telemetry to MQTT.",
            "Agent A detects anomalies → alerts.",
            "Agent B performs RCA → diagnosis.",
            "Agent C creates tickets.",
            "Agent D provides human review and feedback.",
        ], arch_path),
        ("Agent A", add_content_slide, "Agent A: Monitor", [
            "• Subscribes to telemetry/*",
            "• Sliding window (60s/120s) features: mean, std, slope",
            "• Z-score anomaly detection per signal",
            "• Publishes AlertEvent with evidence to alerts/*",
        ]),
        ("Agent B", add_content_slide, "Agent B: Diagnosis", [
            "• Subscribes to alerts, fetches telemetry history",
            "• Rule-based + LLM (ReAct) for root cause analysis",
            "• Outputs: root_cause, confidence, recommended_actions, evidence",
            "• Rules: bearing_wear, clogging, valve_stuck, sensor_drift",
        ]),
        ("Agent C", add_content_slide, "Agent C: Ticket", [
            "• Subscribes to diagnosis",
            "• Creates tickets (GitHub Issues or local API)",
            "• Publishes ticket_id to MQTT",
            "• Queues for Agent D review",
        ]),
        ("Agent D", add_content_slide, "Agent D: Review", [
            "• Human review interface (approve / reject / edit)",
            "• Chat assistant with RAG (similar cases, rules, feedback)",
            "• Optional: Create Salesforce Case on approve",
            "• Scenario management: load, start, stop fault scenarios",
        ]),
        ("Data Flow", add_content_slide, "Data Flow", [
            "Telemetry → Alerts → Diagnosis → Ticket → Feedback",
            "Each step adds structure and context.",
            "Feedback can update rules and thresholds.",
        ], flow_path),
        ("MQTT", add_content_slide, "MQTT Message Bus", [
            "• Lightweight pub/sub for real-time streaming",
            "• Topics: telemetry/{asset}, alerts/{asset}, diagnosis/{asset}, tickets/{asset}",
            "• Decouples agents; each subscribes/publishes independently",
            "• Mosquitto broker (Docker or standalone)",
        ]),
        ("Simulator", add_content_slide, "Simulator & Fault Injection", [
            "• Models pump/piping/bearing subsystem",
            "• Fault types: bearing_wear, clogging, valve_stuck, sensor_drift",
            "• JSON-driven scenarios, reproducible with seed",
            "• Optional: 3D visualization + VLM for multimodal diagnosis",
        ]),
        ("Fault Types", add_content_slide, "Supported Fault Types", [
            "• bearing_wear: Gradual increase in vibration, bearing temp",
            "• clogging: Flow drops, pressure rises, current increases",
            "• valve_stuck: Valve position unchanged, flow unresponsive",
            "• sensor_drift: Single sensor offset, others consistent",
        ]),
        ("Dashboard", add_content_slide, "Agent D Dashboard Features", [
            "• Review Queue: Pending diagnoses, approve/reject",
            "• Alerts: View, regenerate diagnosis, add to queue",
            "• Sensors: Real-time telemetry visualization",
            "• Chat: ReAct assistant with RAG tools",
            "• Scenarios: Load/start/stop simulator scenarios",
        ], dash_path),
        ("Rules", add_content_slide, "Troubleshooting Rules & RAG", [
            "• Rules stored in agent-diagnosis/rules/*.md",
            "• Create from natural language (LLM) or flowchart upload (VLM)",
            "• RAG: query similar diagnoses, alerts, feedback, rules",
            "• Vector index (sqlite-vec + sentence-transformers)",
        ]),
        ("Salesforce", add_content_slide, "Salesforce Integration", [
            "• On approve: optionally create Salesforce Case",
            "• Pre-filled form: Subject, Description, Priority, Status, etc.",
            "• Picklist values fetched from Salesforce org",
            "• TicketConnector interface for extensibility",
        ]),
        ("Tech Stack", add_content_slide, "Technology Stack", [
            "• Backend: Python 3.11, FastAPI, pydantic",
            "• Message bus: MQTT (Mosquitto)",
            "• LLM: LangChain, LangGraph, OpenAI/DeepSeek",
            "• Frontend: React, Vite",
            "• RAG: sqlite-vec, sentence-transformers",
        ]),
        ("Demo Flow", add_content_slide, "Demo Flow", [
            "1. Start MQTT, Simulator, Agents A/B/C/D",
            "2. Load fault scenario (e.g., bearing_wear)",
            "3. Simulator injects fault → Agent A detects → Agent B diagnoses",
            "4. Agent C creates ticket → appears in Review Queue",
            "5. Human approves, optionally creates Salesforce Case",
        ]),
        ("Human Loop", add_content_slide, "Human-in-the-Loop Design", [
            "• Agent D requires human approval before ticket closure",
            "• Approve / Reject / Edit with notes",
            "• Optional: Create Salesforce Case on approve",
            "• Feedback feeds back into rules and RAG for improvement",
        ]),
        ("Benefits", add_content_slide, "Key Benefits", [
            "• End-to-end automated pipeline",
            "• Explainable diagnosis with evidence",
            "• Human-in-the-loop for safety",
            "• Extensible (rules, RAG, Salesforce)",
            "• Reproducible scenarios for evaluation",
        ]),
        ("Evaluation", add_content_slide, "Evaluation & Metrics", [
            "• Detection latency: fault injection → first alert",
            "• Precision/recall: event-level (truth.fault != none)",
            "• Ticket creation latency: diagnosis → ticket",
            "• False positive rate (planned)",
        ]),
        ("Lessons", add_content_slide, "Lessons Learned", [
            "• Multi-agent design enables modularity and scalability",
            "• MQTT simplifies real-time event streaming",
            "• Rules + LLM balances interpretability and flexibility",
            "• Human review is essential for industrial applications",
        ]),
        ("Future Work", add_content_slide, "Future Work", [
            "• Evaluation metrics: precision, recall, latency",
            "• Fine-tuning LLM on diagnosis data",
            "• Multi-asset dashboard enhancements",
            "• Grafana integration for metrics",
        ]),
        ("Conclusion", add_content_slide, "Conclusion", [
            "• Multi-agent system for powerplant monitoring",
            "• Real-time anomaly detection → RCA → ticket → human review",
            "• Combines rules, LLM, RAG, and external integrations",
            "• Practical, extensible, and explainable",
        ]),
        ("Q&A", add_content_slide, "Q & A", [
            "Thank you!",
            "",
            "Questions?",
        ]),
    ]

    for item in slides:
        if len(item) == 4:
            _, fn, a, b = item
            if fn == add_title_slide:
                fn(prs, a, b)
            else:
                fn(prs, a, b, image_path=None)
        else:
            _, fn, a, b, img = item
            fn(prs, a, b, img)

    prs.save(str(output_pptx))
    print(f"Presentation saved to: {output_pptx}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
